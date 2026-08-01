# -*- coding: utf-8 -*-
"""مولّد عروض المحاضرات العملية — معمارية الحاسوب (جامعة إب)"""
import os, sys, math
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from theme_ibbu import *
from theme_ibbu import _rtlcol
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

PRS = None
N = 0
TOTAL = 0

def new():
    global N
    N += 1
    return blank(PRS)

# ============================================================
#  شريحة الغلاف
# ============================================================
def slide_cover(cfg):
    s = dark_slide(PRS)
    # زخارف دائرية
    circle(s, 11.4, -2.6, 6.0, fill=DARK2)
    circle(s, -2.4, 5.4, 4.6, fill=DARK2)
    circle(s, 11.8, 5.2, 3.8, fill=DARK2)
    pill(s, SW - 6.0, 0.5, 5.6, 0.42, "جامعة إب · كلية الحاسابات والعلوم التطبيقية", DARK2, 13, True, "#BFE3C6", HEAD_FONT)
    pill(s, 0.5, 0.5, 2.8, 0.42, f"المحاضرة {cfg['num']}", DARK2, 13, True, "#9FD2A8", HEAD_FONT)
    tf = txbox(s, 1.0, 2.6, SW - 2.0, 1.9)
    para_(tf, cfg["title"], 40, True, WHITE, PP_ALIGN.RIGHT, HEAD_FONT, 0, 0, 1.05, True, True)
    tf2 = txbox(s, 1.0, 4.25, SW - 2.0, 1.0)
    para_(tf2, cfg.get("title_en", ""), 19, False, "#B8E0C4", PP_ALIGN.RIGHT, BODY_FONT, 0, 0, 1.3, True, True)
    # بطاقات بيانات
    meta = [
        ("المقرر", "معمارية الحاسوب — عملي (CS 221)", EMERALD),
        ("المحاضر", cfg.get("instructor", "م. طارق العمري"), SKY),
        ("المدة", cfg.get("duration", "120 دقيقة"), GOLD),
        ("الوحدة", cfg.get("week", ""), "#0F5C37"),
    ]
    x = 0.5
    for t, v, c in meta:
        card(s, x, 5.75, 3.03, 0.95, fill=DARK2, line_color=DARK2, accent=c)
        tf = txbox(s, x + 0.2, 5.87, 2.6, 0.75)
        para_(tf, t, 12.5, True, c, PP_ALIGN.RIGHT, HEAD_FONT, 1, 0, 1.0, True, True)
        para_(tf, v, 11.5, False, "#E5F2E8", PP_ALIGN.RIGHT, BODY_FONT, 0, 0, 1.15, True)
        x += 3.13
    notes(s, cfg.get("cover_note", f"أرحب بالطلاب في {cfg['title']} — المحاضرة العملية {cfg['num']} لمقرر معمارية الحاسوب. أستعرض أهداف المحاضرة ثم أنتقل للجزء العملي على بيئة MARS. أذكر للطلاب تجهيز بيئة العمل قبل البدء."))
    return s

# ============================================================
#  الأهداف
# ============================================================
def slide_objectives(cfg):
    s = new()
    header(s, "الأهداف التعليمية", EMERALD, "أهداف هذه المحاضرة العملية",
           "بنهاية المحاضرة سيكون الطالب قادراً على:", N, TOTAL, cfg["section_label"], badge_w=2.6)
    goals = cfg["objectives"]
    y = 2.3
    for i, g in enumerate(goals):
        col = i % 2
        x = 0.5 + col * 6.35
        if i and i % 2 == 0:
            y += 1.12
        c = SECTION_COLORS[i % len(SECTION_COLORS)]
        card(s, x, y, 6.1, 0.95, accent=c)
        chip(s, x + 0.25, y + 0.27, 0.42, "✓", c, size=15)
        tf = txbox(s, x + 0.85, y + 0.15, 5.05, 0.75)
        para_(tf, g, 15, False, INK, PP_ALIGN.RIGHT, BODY_FONT, 0, 0, 1.2, True, True)
    notes(s, cfg.get("objectives_note", "أقرأ الأهداف وأربطها بأجزاء المحاضرة العملية، وأؤكد أن كل جزء يتضمن تطبيقاً مباشراً على بيئة MARS."))
    return s

# ============================================================
#  مخطط المحاضرة (جدول الوقت)
# ============================================================
def slide_agenda(cfg):
    s = new()
    header(s, "مخطط المحاضرة", GOLD, "توزيع الوقت والأنشطة", None, N, TOTAL, cfg["section_label"], badge_w=2.6)
    rows = cfg["agenda"]
    n = len(rows)
    top, bottom = 2.3, 6.82
    step = (bottom - top) / n
    compact = n > 8
    y = top
    for i, (act, dur, desc) in enumerate(rows):
        c = SECTION_COLORS[i % len(SECTION_COLORS)]
        rh = step - 0.08
        card(s, 0.5, y, 12.33, rh, accent=c)
        cd = 0.32 if compact else 0.38
        chip(s, 0.75, y + (rh - cd) / 2, cd, str(i + 1), c, size=12 if compact else 14)
        tf = txbox(s, 1.35, y + (rh - 0.42) / 2, 8.3, 0.45)
        para_(tf, act, 12.5 if compact else 14.5, True, INK, PP_ALIGN.RIGHT, BODY_FONT, 0, 0, 1.05, True, True)
        if desc:
            para_(tf, "  —  " + desc, 10.5 if compact else 12, False, MUTED, PP_ALIGN.RIGHT, BODY_FONT, 0, 0, 1.05, True)
        pill(s, 12.0, y + (rh - 0.34) / 2, 0.7, 0.34, dur, c, 10.5 if compact else 11.5, True, WHITE, BODY_FONT)
        y += step
    notes(s, cfg.get("agenda_note", "أوضح توزيع وقت المحاضرة: جزء نظري قصير ثم تطبيق عملي مباشر على MARS في كل فقرة."))
    return s

# ============================================================
#  نقاط / قائمة
# ============================================================
def slide_bullets(sl, color):
    s = new()
    header(s, sl["badge"], color, sl["title"], sl.get("subtitle"), N, TOTAL, sl.get("section") or "")
    side = sl.get("side")
    main_w = 8.5 if side else 12.3
    card(s, 0.5, 2.3, main_w, 4.55, accent=color)
    tf = txbox(s, 0.85, 2.55, main_w - 0.7, 4.1)
    first = True
    for it in sl.get("points", []):
        p = para_(tf, "", 0, False, BODY, PP_ALIGN.RIGHT, BODY_FONT, 9, 0, 1.28, True, first)
        first = False
        run(p, "●  ", 13, True, color, BODY_FONT, icon=True)
        if isinstance(it, tuple):
            run(p, it[0], 15, True, INK, BODY_FONT)
            run(p, " — " + it[1], 15, False, BODY, BODY_FONT)
        else:
            run(p, it, 15, False, BODY, BODY_FONT)
    if side:
        card(s, 9.2, 2.3, 3.63, 4.55, fill="F0F7F1", line_color="#C9E3CD", accent=GOLD)
        tf = txbox(s, 9.45, 2.55, 3.15, 0.5)
        para_(tf, side.get("title", "نقطة محورية"), 15, True, DARK, PP_ALIGN.RIGHT, HEAD_FONT, 0, 0, 1, True, True)
        tf = txbox(s, 9.45, 3.1, 3.15, 3.6)
        first = True
        for it in side.get("points", []):
            p = para_(tf, "", 0, False, BODY, PP_ALIGN.RIGHT, BODY_FONT, 8, 0, 1.25, True, first)
            first = False
            run(p, "★  ", 12, True, GOLD, BODY_FONT, icon=True)
            run(p, it, 13.5, False, BODY, BODY_FONT)
    notes(s, sl.get("note", f"أشرح نقاط «{sl['title']}» بالتفصيل مع أمثلة من بيئة MARS."))
    return s

# ============================================================
#  خطوات مرقمة
# ============================================================
def slide_steps(sl, color):
    s = new()
    header(s, sl["badge"], color, sl["title"], sl.get("subtitle"), N, TOTAL, sl.get("section") or "")
    steps = sl["steps"]
    y = 2.3
    for i, st in enumerate(steps):
        card(s, 0.5, y, 12.33, 0.68, accent=color)
        chip(s, 0.75, y + 0.13, 0.42, str(i + 1), color, size=15)
        tf = txbox(s, 1.4, y + 0.08, 11.2, 0.55)
        para_(tf, st, 15, False, INK, PP_ALIGN.RIGHT, BODY_FONT, 0, 0, 1.15, True, True)
        y += 0.78
    if sl.get("note_line"):
        tf = txbox(s, 0.5, y + 0.05, 12.3, 0.5)
        para_(tf, sl["note_line"], 13, True, DARK, PP_ALIGN.RIGHT, BODY_FONT, 0, 0, 1.2, True, True)
    notes(s, sl.get("note", f"أنفذ الخطوات {len(steps)} أمام الطلاب على الشاشة، وأطلب منهم التكرار خطوة بخطوة."))
    return s

# ============================================================
#  مخطط / صورة
# ============================================================
def slide_diagram(sl, color):
    s = new()
    header(s, sl["badge"], color, sl["title"], sl.get("subtitle"), N, TOTAL, sl.get("section") or "")
    path = sl["img"]
    from PIL import Image
    iw, ih = Image.open(path).size
    ratio = iw / ih
    side = sl.get("side")
    if ratio < 1.1:
        # صورة عمود يمين + نقاط يسار
        w, h = diagram(s, path, 7.5, 2.5, 5.1, 4.05, caption=sl.get("caption"))
        card(s, 0.5, 2.3, 6.7, 4.55, accent=color)
        tf = txbox(s, 0.85, 2.55, 6.0, 4.1)
        first = True
        for it in (side or sl.get("points", [])):
            p = para_(tf, "", 0, False, BODY, PP_ALIGN.RIGHT, BODY_FONT, 9, 0, 1.3, True, first)
            first = False
            run(p, "●  ", 13, True, color, BODY_FONT, icon=True)
            if isinstance(it, tuple):
                run(p, it[0], 15, True, INK, BODY_FONT)
                run(p, " — " + it[1], 15, False, BODY, BODY_FONT)
            else:
                run(p, it, 15, False, BODY, BODY_FONT)
    else:
        # صورة أفقية أعلى + نقاط أسفل
        cap_h = 0.3 if sl.get("caption") else 0.0
        pts = sl.get("points", [])
        img_h = 4.15 if not pts else 3.1
        w, h = diagram(s, path, None, 2.35, 12.0, img_h, caption=sl.get("caption"))
        if pts:
            y = 5.65 if not sl.get("caption") else 6.05
            card(s, 0.5, y, 12.33, 1.0, fill="F0F7F1", line_color="#C9E3CD", accent=GOLD)
            tf = txbox(s, 0.85, y + 0.12, 11.7, 0.8)
            first = True
            for it in pts:
                p = para_(tf, "", 0, False, BODY, PP_ALIGN.RIGHT, BODY_FONT, 3, 0, 1.2, True, first)
                first = False
                run(p, "★  ", 12, True, GOLD, BODY_FONT, icon=True)
                run(p, it, 14, False, BODY, BODY_FONT)
    notes(s, sl.get("note", f"أشرح مخطط «{sl['title']}» وأربطه بالمفاهيم المطروحة."))
    return s

# ============================================================
#  كود + شرح
# ============================================================
def slide_code(sl, color):
    s = new()
    header(s, sl["badge"], color, sl["title"], sl.get("subtitle"), N, TOTAL, sl.get("section") or "")
    lines = sl["code"] if isinstance(sl["code"], list) else sl["code"].split("\n")
    explain = sl.get("explain", [])
    w = 7.3 if explain else 12.3
    # تحديد عدد الأسطر المتسق مع ارتفاع الصندوق
    size = sl.get("size", 12)
    lh = size * 1.38 / 72.0
    box_h = 4.55
    title_pad = 0.5 if sl.get("file") else 0.3
    max_lines = int((box_h - title_pad) / lh)
    truncated = len(lines) > max_lines
    shown = lines[: max_lines - 1 if truncated else max_lines]
    if truncated:
        shown = shown + ["#  … بقية السطور في ملف الـ .asm الكامل"]
    code_box(s, 0.5, 2.3, w, box_h, shown, size=size, title=sl.get("file"))
    if explain:
        card(s, 8.0, 2.3, 4.83, 4.55, accent=color)
        tf = txbox(s, 8.25, 2.5, 4.35, 4.2)
        first = True
        for it in explain:
            p = para_(tf, "", 0, False, BODY, PP_ALIGN.RIGHT, BODY_FONT, 8, 0, 1.25, True, first)
            first = False
            run(p, "✓  ", 13, True, color, BODY_FONT, icon=True)
            if isinstance(it, tuple):
                run(p, it[0], 14, True, INK, BODY_FONT)
                run(p, " — " + it[1], 14, False, BODY, BODY_FONT)
            else:
                run(p, it, 14, False, BODY, BODY_FONT)
    notes(s, sl.get("note", f"أشرح سطراً بسطر في «{sl['title']}»، وألفت النظر لترتيب syscall والتسميات."))
    return s

# ============================================================
#  جدول
# ============================================================
def data_table(slide, x, y, w, headers, rows, col_w=None, size=13, header_fill=DARK,
               mono_cols=(), col_align=None):
    nrows = len(rows) + 1
    ncols = len(headers)
    gfx = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(0.4))
    tbl = gfx.table
    if col_w:
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Inches(cw)
    # تقدير ارتفاع كل صف
    def est_lines(txt, col_i, sz):
        cw = col_w[col_i] if col_w else w / ncols
        cpl = max(1, int((cw - 0.16) / (0.5 * sz / 72.0)))
        return max(1, math.ceil(len(txt) / cpl))
    for r in range(nrows):
        maxl = 1
        for c in range(ncols):
            txt = headers[c] if r == 0 else str(rows[r - 1][c])
            maxl = max(maxl, est_lines(txt, c, size))
        tbl.rows[r].height = Inches(0.16 + maxl * 0.21)
    for r in range(nrows):
        for c in range(ncols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = C(header_fill)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = C("F0F7F1" if r % 2 == 0 else CARD)
            tf = cell.text_frame
            tf.word_wrap = True
            _rtlcol(tf)
            text = headers[c] if r == 0 else str(rows[r - 1][c])
            if c in mono_cols:
                align = PP_ALIGN.LEFT
                font = MONO_FONT
                rtl = False
            else:
                align = col_align[c] if col_align else (PP_ALIGN.RIGHT if c == 0 else PP_ALIGN.CENTER)
                font = HEAD_FONT if r == 0 else BODY_FONT
                rtl = True
            p = para_(tf, "", size=(size if r > 0 else size), bold=(r == 0),
                      color=(WHITE if r == 0 else BODY), align=align, font=font,
                      space_after=0, line=1.1, rtl=rtl, first=True)
            run(p, text, size if r > 0 else size, r == 0,
                WHITE if r == 0 else (INK if c == 0 else BODY), font)
    return gfx

def slide_table(sl, color):
    s = new()
    header(s, sl["badge"], color, sl["title"], sl.get("subtitle"), N, TOTAL, sl.get("section") or "")
    headers = sl["headers"]
    rows = sl["rows"]
    nrows = len(rows) + 1
    data_table(s, 0.5, 2.35, 12.33, headers, rows, col_w=sl.get("col_w"),
               size=sl.get("size", 13), mono_cols=sl.get("mono_cols", ()), col_align=sl.get("col_align"))
    if sl.get("note_line"):
        tf = txbox(s, 0.5, 6.55, 12.3, 0.5)
        para_(tf, sl["note_line"], 13.5, True, DARK, PP_ALIGN.RIGHT, BODY_FONT, 0, 0, 1.2, True, True)
    notes(s, sl.get("note", f"أشرح جدول «{sl['title']}» صفاً صفاً وأوضح العلاقات بين العناصر."))
    return s

# ============================================================
#  نشاط عملي
# ============================================================
def slide_activity(act, cfg):
    s = new()
    header(s, "نشاط عملي", GOLD, act["title"], act.get("subtitle"), N, TOTAL, cfg["section_label"], badge_w=2.4)
    steps = act["steps"]
    y = 2.3
    for i, st in enumerate(steps):
        card(s, 0.5, y, 12.33, 0.68, accent=GOLD)
        chip(s, 0.75, y + 0.13, 0.42, str(i + 1), GOLD, size=15)
        tf = txbox(s, 1.4, y + 0.08, 11.2, 0.55)
        para_(tf, st, 15, False, INK, PP_ALIGN.RIGHT, BODY_FONT, 0, 0, 1.15, True, True)
        y += 0.78
    if act.get("note_line"):
        tf = txbox(s, 0.5, y + 0.05, 12.3, 0.5)
        para_(tf, act["note_line"], 13, True, DARK, PP_ALIGN.RIGHT, BODY_FONT, 0, 0, 1.2, True, True)
    notes(s, act.get("note", "أمنح الطلاب وقتاً لتطبيق النشاط على أجهزتهم، وأتجول لمساعدتهم، ثم أناقش النتائج مع الجميع."))
    return s

# ============================================================
#  واجب منزلي
# ============================================================
def slide_homework(hw, cfg):
    s = new()
    header(s, "الواجب المنزلي", SKY, hw["title"], hw.get("subtitle"), N, TOTAL, cfg["section_label"], badge_w=2.4)
    items = hw["items"]
    card(s, 0.5, 2.3, 7.4, 4.55, accent=SKY)
    tf = txbox(s, 0.85, 2.55, 6.7, 4.1)
    first = True
    for it in items:
        p = para_(tf, "", 0, False, BODY, PP_ALIGN.RIGHT, BODY_FONT, 9, 0, 1.28, True, first)
        first = False
        run(p, "✓  ", 13, True, SKY, BODY_FONT, icon=True)
        run(p, it, 14.5, False, BODY, BODY_FONT)
    cr = hw.get("criteria", [])
    if cr:
        card(s, 8.1, 2.3, 4.73, 4.55, fill="F0F7FF", line_color="#C5E3F7", accent=GOLD)
        tf = txbox(s, 8.35, 2.5, 4.25, 0.5)
        para_(tf, hw.get("criteria_title", "معايير التسليم"), 15, True, DARK, PP_ALIGN.RIGHT, HEAD_FONT, 0, 0, 1, True, True)
        tf = txbox(s, 8.35, 3.1, 4.25, 3.6)
        first = True
        for it in cr:
            p = para_(tf, "", 0, False, BODY, PP_ALIGN.RIGHT, BODY_FONT, 8, 0, 1.25, True, first)
            first = False
            run(p, "★  ", 12, True, GOLD, BODY_FONT, icon=True)
            run(p, it, 13, False, BODY, BODY_FONT)
    notes(s, hw.get("note", "أشرح الواجب بوضوح وأذكر موعد التسليم وطريقة رفعه."))
    return s

# ============================================================
#  ملخص
# ============================================================
def slide_summary(sum_, cfg):
    s = new()
    header(s, "الملخص", DARK, "الملخص — أهم ما تعلمناه اليوم", None, N, TOTAL, cfg["section_label"], badge_w=2.4)
    y = 2.3
    for i, it in enumerate(sum_):
        col = i % 2
        x = 0.5 + col * 6.35
        if i and i % 2 == 0:
            y += 1.1
        c = SECTION_COLORS[i % len(SECTION_COLORS)]
        card(s, x, y, 6.1, 0.95, accent=c)
        chip(s, x + 0.25, y + 0.27, 0.42, str(i + 1), c, size=15)
        tf = txbox(s, x + 0.85, y + 0.15, 5.05, 0.75)
        para_(tf, it, 15, False, INK, PP_ALIGN.RIGHT, BODY_FONT, 0, 0, 1.2, True, True)
    notes(s, sum_ and "ألخص أبرز النقاط وأفتح باب الأسئلة قبل الختام." or "")
    return s

# ============================================================
#  ختام
# ============================================================
def slide_closing(cfg):
    s = dark_slide(PRS)
    circle(s, 11.6, -2.6, 5.8, fill=DARK2)
    circle(s, -2.2, 5.3, 4.4, fill=DARK2)
    pill(s, SW - 4.9, 0.6, 4.4, 0.4, cfg["section_label"], DARK2, 12.5, True, "#BFE3C6", HEAD_FONT)
    tf = txbox(s, 1.0, 3.05, SW - 2.0, 1.5)
    para_(tf, "شكراً لحسن استماعكم ومشاركتكم", 38, True, WHITE, PP_ALIGN.RIGHT, HEAD_FONT, 0, 0, 1.1, True, True)
    tf2 = txbox(s, 1.0, 4.35, SW - 2.0, 0.8)
    para_(tf2, cfg.get("closing_note", "أتمنى أن تكون المحاضرة مفيدة، وإلى اللقاء في المحاضرة القادمة."),
          17, False, "#B8E0C4", PP_ALIGN.RIGHT, BODY_FONT, 0, 0, 1.3, True, True)
    if cfg.get("next"):
        pill(s, 1.0, 5.6, 5.2, 0.55, "المحاضرة القادمة: " + cfg["next"], "#0B3A22", 14, True, "#9FD2A8", HEAD_FONT)
    notes(s, "أشكر الطلاب وأذكّرهم بالواجب والمحاضرة القادمة.")
    return s

# ============================================================
#  المُجمّع الرئيسي
# ============================================================
def build_lecture(cfg, out_path):
    global PRS, N, TOTAL
    PRS = new_presentation()
    N = 1  # الغلاف يُعدّ شريحة 1 (بلا رقم في التذييل)
    slides = cfg.get("slides", [])
    new_calls = 1 + (1 if cfg.get("agenda") else 0) + len(slides)
    for k in ("activity", "homework", "summary"):
        if cfg.get(k):
            new_calls += 1
    TOTAL = new_calls + 2  # الغلاف + الختام

    slide_cover(cfg)
    slide_objectives(cfg)
    if cfg.get("agenda"):
        slide_agenda(cfg)
    for si, sl in enumerate(slides):
        color = sl.get("color") or SECTION_COLORS[si % len(SECTION_COLORS)]
        st = sl["type"]
        if st == "bullets":
            slide_bullets(sl, color)
        elif st == "steps":
            slide_steps(sl, color)
        elif st == "diagram":
            slide_diagram(sl, color)
        elif st == "code":
            slide_code(sl, color)
        elif st == "table":
            slide_table(sl, color)
        else:
            slide_bullets(sl, color)
    if cfg.get("activity"):
        slide_activity(cfg["activity"], cfg)
    if cfg.get("homework"):
        slide_homework(cfg["homework"], cfg)
    if cfg.get("summary"):
        slide_summary(cfg["summary"], cfg)
    slide_closing(cfg)

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    PRS.save(out_path)
    print("SAVED:", out_path, "| slides:", len(PRS.slides._sldIdLst))
    return len(PRS.slides._sldIdLst)
