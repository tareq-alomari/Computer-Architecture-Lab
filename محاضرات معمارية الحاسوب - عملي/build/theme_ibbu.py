# -*- coding: utf-8 -*-
"""نظام تصميم موحّد لعروض معمارية الحاسوب — جامعة إب"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)
SW, SH = 13.333, 7.5

# ============ هوية جامعة إب ============
DARK    = "0C4A2D"   # Royal Green (داكن)
DARK2   = "0F5C37"
EMERALD = "15803D"   # Emerald
GREEN   = "059669"   # أخضر ناصع
SKY     = "0EA5E9"   # Sky Blue
SKY2    = "0284C7"
GOLD    = "C79100"   # Gold
GOLD2   = "A16207"
NAVY    = "1A1A2E"   # Deep Navy
INK     = "1A1A2E"
BODY    = "3F4756"
MUTED   = "6B7280"
BG      = "F5F8F5"   # خلفية فاتحة مخضرة
CARD    = "FFFFFF"
LINE    = "DDE6DD"
WHITE   = "FFFFFF"
RED     = "DC2626"
VIOLET  = "7C3AED"

# ألوان الأقسام (تُدار تلقائياً عند كل جزء)
SECTION_COLORS = [EMERALD, SKY, GOLD, DARK, "0F5C37", "0284C7", "A16207", "059669"]

HEAD_FONT = "Noto Kufi Arabic"
BODY_FONT = "Noto Sans Arabic"
MONO_FONT = "DejaVu Sans Mono"
ICON_FONT = "DejaVu Sans"

def C(hexstr):
    return RGBColor.from_string(hexstr.lstrip("#"))

def _set_fonts(run, name):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)

def run(para, text, size=18, bold=False, color=BODY, font=BODY_FONT, icon=False):
    r = para.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = C(color)
    _set_fonts(r, ICON_FONT if icon else font)
    return r

def para_(tf, text="", size=18, bold=False, color=BODY, align=PP_ALIGN.RIGHT,
          font=BODY_FONT, space_after=6, space_before=0, line=1.12, rtl=True, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line
    if rtl:
        pPr = p._p.get_or_add_pPr()
        pPr.set("rtl", "1")
    if text:
        run(p, text, size, bold, color, font)
    return p

def _rtlcol(tf):
    bp = tf._txBody.find(qn("a:bodyPr"))
    if bp is not None:
        bp.set("rtlCol", "1")

def txbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    _rtlcol(tf)
    return tf

def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = C(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = C(line)
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp

def circle(slide, x, y, d, fill=None, line=None):
    return rect(slide, x, y, d, d, fill, line, shape=MSO_SHAPE.OVAL)

def line(slide, x, y, w, h=0.0, color=LINE, width=1.0):
    return rect(slide, x, y, w, h if h else 0.014, fill=color)

def full_bg(slide, color=BG):
    rect(slide, -0.02, -0.02, SW + 0.05, SH + 0.05, fill=color)

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def text_in_shape(sp, text, size=16, bold=True, color=WHITE, font=BODY_FONT, align=PP_ALIGN.CENTER):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    _rtlcol(tf)
    p = para_(tf, text, size, bold, color, align, font, 0, 0, 1.0, True, True)
    pPr = p._p.get_or_add_pPr()
    pPr.set("rtl", "1")
    return sp

def pill(slide, x, y, w, h, text, fill, size=13, bold=True, color=WHITE, font=HEAD_FONT):
    sp = rect(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    text_in_shape(sp, text, size, bold, color, font)
    return sp

def chip(slide, x, y, d, symbol, fill, fg=WHITE, size=16):
    c = circle(slide, x, y, d, fill=fill)
    text_in_shape(c, symbol, size, True, fg, ICON_FONT)
    return c

def card(slide, x, y, w, h, fill=CARD, line_color=LINE, accent=None, accent_h=0.06):
    sp = rect(slide, x, y, w, h, fill=fill, line=line_color, line_w=1.0,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.055)
    if accent:
        rect(slide, x + 0.02, y + 0.02, w - 0.04, accent_h, fill=accent)
    return sp

def header(slide, badge, badge_color, title, subtitle=None, num=None, total=None, section="",
           badge_w=None):
    full_bg(slide)
    rect(slide, 0, 0, SW, 0.07, fill=badge_color)
    rect(slide, 0, 0, 0.09, SH, fill=DARK)
    bw = badge_w if badge_w else (len(badge) * 0.135 + 0.42)
    pill(slide, SW - bw - 0.42, 0.40, bw, 0.34, badge, badge_color, 12)
    tf = txbox(slide, 0.8, 0.95, SW - 2.3, 0.9)
    para_(tf, title, 27, True, DARK, PP_ALIGN.RIGHT, HEAD_FONT, 0, 0, 1.05, True, True)
    rect(slide, SW - 1.35, 1.62, 0.55, 0.045, fill=badge_color)
    if subtitle:
        tf2 = txbox(slide, 0.8, 1.72, SW - 2.3, 0.5)
        para_(tf2, subtitle, 14.5, False, MUTED, PP_ALIGN.RIGHT, BODY_FONT, 0, 0, 1.1, True, True)
    footer(slide, num, total, section)

def footer(slide, num=None, total=None, section=""):
    rect(slide, 0, SH - 0.38, SW, 0.012, fill=LINE)
    if section:
        tf = txbox(slide, 0.35, SH - 0.33, 6.0, 0.3)
        para_(tf, section, 10.5, False, MUTED, PP_ALIGN.LEFT, BODY_FONT, 0, 0, 1, False, True)
    if num is not None:
        tf = txbox(slide, SW - 1.6, SH - 0.33, 1.2, 0.3)
        para_(tf, f"{num} / {total}", 11, True, DARK, PP_ALIGN.RIGHT, BODY_FONT, 0, 0, 1, True, True)

def bullet_list(slide, x, y, w, items, size=15, gap=8, marker_color=EMERALD, text_color=BODY,
                marker="●", line=1.22, max_h=None):
    tf = txbox(slide, x, y, w, max_h or 3.0)
    first = True
    for it in items:
        p = para_(tf, "", 0, False, text_color, PP_ALIGN.RIGHT, BODY_FONT,
                  gap, 0, line, True, first)
        first = False
        if isinstance(it, tuple):
            head, body = it
            run(p, marker + "  ", size - 2, True, marker_color, BODY_FONT, icon=True)
            run(p, head, size, True, INK, BODY_FONT)
            run(p, " — " + body, size, False, text_color, BODY_FONT)
        else:
            run(p, marker + "  ", size - 2, True, marker_color, BODY_FONT, icon=True)
            run(p, it, size, False, text_color, BODY_FONT)
    return tf

def add_notes(slide, text):
    notes(slide, text)

def new_presentation():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def picture(slide, path, x, y, w=None, h=None):
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    return slide.shapes.add_picture(path, Inches(x), Inches(y), **kw)

def dark_slide(prs, gradient_top=DARK, gradient_bottom=DARK2):
    s = blank(prs)
    rect(s, -0.02, -0.02, SW + 0.05, SH + 0.05, fill=gradient_top)
    return s

# ================= كتلة الكود (MIPS) — اتجاه LTR =================
CODE_BG  = "0B1E17"   # خلفية داكنة للكود
CODE_LN  = "0E2A20"
TOK_COMMENT = "5E8B6A"
TOK_DIRECT  = "E8B66C"   # ذهبي
TOK_LABEL   = "7FD1F5"   # سماوي
TOK_REG     = "7BE3A1"   # أخضر فاتح
TOK_STR     = "F2A8C6"   # وردي
TOK_NUM     = "C6A4F5"   # بنفسجي
TOK_TEXT    = "E8EDE9"   # أبيض مخضر

def _code_runs(p, text, size, base=TOK_TEXT):
    """تلوين بسيط لأسطر MIPS: تعليق / توجيه / تسمية / مسجل / نص"""
    def emit(t, c, bold=False):
        run(p, t, size, bold, c, MONO_FONT)
    stripped = text.lstrip()
    pad = text[:len(text) - len(stripped)]
    if pad:
        emit(pad, "17332A")
    if stripped.startswith("#"):
        emit(stripped, TOK_COMMENT)
        return
    # فصل التعليق
    code = stripped
    comment = ""
    idx = code.find("#")
    if idx != -1:
        code, comment = code[:idx], code[idx:]
    # تسمية
    if ":" in code and not code.strip().startswith("."):
        parts = code.split(":", 1)
        emit(parts[0].strip() + ":", TOK_LABEL, True)
        code = parts[1]
    # تلوين الرموز
    import re
    pos = 0
    pattern = re.compile(r"(\.\w+|\\\$[a-z0-9]+|\\\$f\w+|\"[^\"]*\"|-?\b\d+\b|\b\w+\b)")
    for m in pattern.finditer(code):
        tok = m.group(0)
        if m.start() > pos:
            emit(code[pos:m.start()], base)
        if tok.startswith("."):
            emit(tok, TOK_DIRECT, True)
        elif tok.startswith("$"):
            emit(tok, TOK_REG)
        elif tok.startswith('"'):
            emit(tok, TOK_STR)
        elif tok.lstrip("-").isdigit():
            emit(tok, TOK_NUM)
        else:
            emit(tok, TOK_TEXT)
        pos = m.end()
    if pos < len(code):
        emit(code[pos:], base)
    if comment:
        emit(" " + comment, TOK_COMMENT)

def code_box(slide, x, y, w, h, code_lines, size=12.5, title=None, title_color=SKY,
             max_lines=None):
    """صندوق كود داكن باتجاه LTR مع تلوين ورقم السطر"""
    box = rect(slide, x, y, w, h, fill=CODE_BG, line="153326", line_w=1.0,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.035)
    title_pad = 0.46 if title else 0.18
    if max_lines is None:
        lh = size * 1.38 / 72.0
        max_lines = int((h - title_pad) / lh)
    lines = code_lines[:max_lines]
    start_y = y + title_pad
    tf = txbox(slide, x + 0.25, start_y, w - 1.35, h - 0.3)
    first = True
    for ln in lines:
        p = para_(tf, "", 0, False, TOK_TEXT, PP_ALIGN.LEFT, MONO_FONT, 0, 0, 1.0, False, first)
        first = False
        _code_runs(p, ln.rstrip("\n"), size)
    # أرقام الأسطر (يسار)
    if len(lines) <= 40:
        tf2 = txbox(slide, x + w - 0.55, start_y, 0.45, h - 0.3)
        first = True
        for i, ln in enumerate(lines, 1):
            p = para_(tf2, "", 0, False, TOK_COMMENT, PP_ALIGN.LEFT, MONO_FONT, 0, 0, 1.0, False, first)
            first = False
            run(p, str(i), size, False, TOK_COMMENT, MONO_FONT)
    if title:
        tf3 = txbox(slide, x + 0.3, y + 0.06, w - 0.6, 0.3)
        para_(tf3, title, 12, True, title_color, PP_ALIGN.LEFT, MONO_FONT, 0, 0, 1, False, True)
    return box

# ================= صورة مخطط (تتلاءم مع النسبة) =================
def diagram(slide, path, x, y, max_w, max_h, caption=None, frame=True, bg=WHITE):
    from PIL import Image
    iw, ih = Image.open(path).size
    ratio = iw / ih
    w, h = max_w, max_w / ratio
    if h > max_h:
        h = max_h
        w = h * ratio
    if x is None:
        x = (SW - w) / 2
    if frame:
        card(slide, x - 0.12, y - 0.12, w + 0.24, h + 0.24, fill=bg, line_color=LINE)
    picture(slide, path, x, y, w=w, h=h)
    if caption:
        tf = txbox(slide, x - 0.12, y + h + 0.06, w + 0.24, 0.3)
        para_(tf, caption, 11, True, MUTED, PP_ALIGN.CENTER, BODY_FONT, 0, 0, 1, True, True)
    return w, h

def step_chip(slide, x, y, d, num, color):
    c = circle(slide, x, y, d, fill=color)
    text_in_shape(c, str(num), 15, True, WHITE, HEAD_FONT)
    return c
